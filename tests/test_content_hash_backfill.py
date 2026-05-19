import shutil
from pathlib import Path

import pytest
from alembic.config import Config
from pptx import Presentation
from sqlalchemy import create_engine, text
from sqlalchemy.orm import close_all_sessions
from sqlalchemy.pool import NullPool

from alembic import command
from bookworm.annotation import NoteTaker
from bookworm.database import init_database
from bookworm.database.models import (
    Book,
    Bookmark,
    DocumentPositionInfo,
    Note,
    PinnedDocument,
    Quote,
    RecentDocument,
)
from bookworm.document import BaseDocument, BasePage, BookMetadata, Pager, Section, create_document
from bookworm.document.uri import DocumentUri
from bookworm.gui.book_viewer import recents_manager
from bookworm.structured_text import (
    CURRENT_CONTENT_HASH_VERSION,
    CURRENT_POSITION_MODEL_VERSION,
    SemanticElementType,
    TEXT_OBJECT_REPLACEMENT_CHAR,
)


def _make_alembic_config(db_url):
    cfg = Config(Path("alembic.ini"))
    cfg.attributes["configure_logger"] = False
    cfg.set_main_option("script_location", "alembic")
    cfg.set_main_option("sqlalchemy.url", db_url)
    return cfg


def _make_textless_presentation(filename):
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(filename)


class _SyntheticPaginatedPage(BasePage):
    def get_text(self):
        return self.document.page_texts[self.index]


class _SyntheticPaginatedDocument(BaseDocument):
    __internal__ = True
    format = "test_paginated_hash"
    extensions = ()

    def __init__(self, uri, page_texts, *, title="Paged Document"):
        super().__init__(uri)
        self.page_texts = tuple(page_texts)
        self._metadata = BookMetadata(title=title, author="")

    def __len__(self):
        return len(self.page_texts)

    def read(self):
        super().read()

    def get_page(self, index):
        return _SyntheticPaginatedPage(self, index)

    @property
    def toc_tree(self):
        return Section(title=self.metadata.title, pager=Pager(first=0, last=len(self) - 1))

    @property
    def metadata(self):
        return self._metadata


def _make_synthetic_paginated_document(path, *page_texts, title="Paged Document"):
    document = _SyntheticPaginatedDocument(
        DocumentUri(
            format=_SyntheticPaginatedDocument.format,
            path=path,
            openner_args={},
        ),
        page_texts,
        title=title,
    )
    document.read()
    return document


def test_loading_existing_records_backfills_hashes_and_preserves_annotations(
    asset, reader, tmp_path
):
    original_uri = DocumentUri.from_filename(asset("roman.epub"))
    original_doc = create_document(original_uri)
    content_hash = original_doc.get_content_hash()

    existing_book = Book.get_or_create(title=original_doc.metadata.title, uri=original_uri)
    existing_doc_info = DocumentPositionInfo.get_or_create(
        title=original_doc.metadata.title,
        uri=original_uri,
    )
    Note.session().add(
        Note(
            title="test",
            content="test note",
            page_number=0,
            position=0,
            section_title="test section",
            section_identifier="test-section",
            book_id=existing_book.id,
        )
    )
    Note.session().commit()

    reader.load(original_uri)

    updated_book = Book.query.one()
    updated_doc_info = DocumentPositionInfo.query.one()
    assert updated_book.id == existing_book.id
    assert updated_doc_info.id == existing_doc_info.id
    assert updated_book.content_hash == content_hash
    assert updated_doc_info.content_hash == content_hash
    reader.unload()

    moved_path = shutil.copy(Path(original_uri.path), tmp_path / "roman.epub")
    moved_uri = DocumentUri.from_filename(moved_path)
    reader.load(moved_uri)

    assert Book.query.count() == 1
    assert DocumentPositionInfo.query.count() == 1
    assert Book.query.one().uri == moved_uri
    assert DocumentPositionInfo.query.one().uri == moved_uri
    assert NoteTaker(reader).get_for_page(0).count() == 1


def test_legacy_image_placeholder_positions_are_migrated_on_load(reader, tmp_path):
    html_path = tmp_path / "image.html"
    html_path.write_text(
        """
        <html>
            <head><title>Book</title></head>
            <body>
                <p>before <img src="pic.png" alt="Chart"> middle
                <img src="pic.png" alt="Chart 2"> omega text</p>
            </body>
        </html>
        """,
        encoding="utf-8",
    )
    (tmp_path / "pic.png").write_bytes(b"same-image-bytes")
    uri = DocumentUri.from_filename(html_path)
    document = create_document(uri)
    try:
        legacy_text = document.get_legacy_content()
        storage_text = document.get_storage_content()
        legacy_hash = document.get_legacy_content_hash()
        legacy_target = legacy_text.index("omega")
        storage_target = storage_text.index("omega")
        storage_selection = slice(storage_target, storage_target + len("omega"))
        section = document.get_page(0).section

        session = Book.session()
        book = Book(
            title=document.metadata.title,
            uri=uri,
            content_hash=legacy_hash,
            content_hash_version=None,
        )
        session.add(book)
        session.flush()
        session.add(
            DocumentPositionInfo(
                title=document.metadata.title,
                uri=uri,
                content_hash=legacy_hash,
                content_hash_version=None,
                last_page=0,
                last_position=legacy_target,
                position_version=None,
            )
        )
        session.add(
            Bookmark(
                title="bookmark",
                page_number=0,
                position=legacy_target,
                section_title=section.title,
                section_identifier=section.unique_identifier,
                book_id=book.id,
                position_version=None,
            )
        )
        session.add(
            Note(
                title="note",
                content="note content",
                page_number=0,
                position=legacy_target,
                start_pos=legacy_target,
                end_pos=legacy_target + len("omega"),
                section_title=section.title,
                section_identifier=section.unique_identifier,
                book_id=book.id,
                position_version=None,
            )
        )
        session.add(
            Quote(
                title="quote",
                content="quote content",
                page_number=0,
                position=legacy_target,
                start_pos=legacy_target,
                end_pos=legacy_target + len("omega"),
                section_title=section.title,
                section_identifier=section.unique_identifier,
                book_id=book.id,
                position_version=None,
            )
        )
        session.commit()
        session.execute(text("UPDATE book SET content_hash_version = NULL"))
        session.execute(
            text(
                "UPDATE document_position_info "
                "SET content_hash_version = NULL, position_version = NULL"
            )
        )
        for table_name in ("bookmark", "note", "quote"):
            session.execute(text(f"UPDATE {table_name} SET position_version = NULL"))
        session.commit()
        session.expire_all()

        reader.load(uri)

        assert Book.query.one().content_hash_version == CURRENT_CONTENT_HASH_VERSION
        assert DocumentPositionInfo.query.one().position_version == CURRENT_POSITION_MODEL_VERSION
        assert Bookmark.query.one().position_version == CURRENT_POSITION_MODEL_VERSION
        assert Note.query.one().position_version == CURRENT_POSITION_MODEL_VERSION
        assert Quote.query.one().position_version == CURRENT_POSITION_MODEL_VERSION
        assert DocumentPositionInfo.query.one().last_position == storage_target
        assert Bookmark.query.one().position == storage_target
        assert Note.query.one().position == storage_target
        assert Quote.query.one().position == storage_target
        assert Note.query.one().start_pos == storage_target
        assert Quote.query.one().start_pos == storage_target
        assert Note.query.one().end_pos == storage_selection.stop
        assert Quote.query.one().end_pos == storage_selection.stop

        reader.view.set_insertion_point(document.get_content().index("omega") + 1)
        reader.save_current_position()
        assert (
            DocumentPositionInfo.query.one().last_position
            == document.display_to_storage_position(document.get_content().index("omega") + 1)
        )
    finally:
        document.close()
        reader.unload()


def test_current_records_merge_legacy_duplicates_on_load(reader, tmp_path):
    original_path = tmp_path / "original.html"
    moved_path = tmp_path / "moved.html"
    html_content = """
    <html>
        <head><title>Book</title></head>
        <body>
            <p>before <img src="pic.png" alt="Chart"> middle
            <img src="pic.png" alt="Chart 2"> after text</p>
        </body>
    </html>
    """
    (tmp_path / "pic.png").write_bytes(b"same-image-bytes")
    original_path.write_text(html_content, encoding="utf-8")
    shutil.copy(original_path, moved_path)

    original_uri = DocumentUri.from_filename(original_path)
    moved_uri = DocumentUri.from_filename(moved_path)
    document = create_document(moved_uri)
    try:
        current_hash = document.get_content_hash()
        legacy_hash = document.get_legacy_content_hash()
        legacy_position = document.get_legacy_content().index("after text")
        storage_position = document.get_storage_content().index("after text")
        section = document.get_page(0).section

        session = Book.session()
        current_book = Book(
            title=document.metadata.title,
            uri=moved_uri,
            content_hash=current_hash,
            content_hash_version=CURRENT_CONTENT_HASH_VERSION,
        )
        legacy_book = Book(
            title=document.metadata.title,
            uri=original_uri,
            content_hash=legacy_hash,
            content_hash_version=None,
        )
        session.add_all((current_book, legacy_book))
        session.flush()
        current_position_info = DocumentPositionInfo(
            title=document.metadata.title,
            uri=moved_uri,
            content_hash=current_hash,
            content_hash_version=CURRENT_CONTENT_HASH_VERSION,
            last_page=0,
            last_position=0,
            position_version=CURRENT_POSITION_MODEL_VERSION,
        )
        legacy_position_info = DocumentPositionInfo(
            title=document.metadata.title,
            uri=original_uri,
            content_hash=legacy_hash,
            content_hash_version=None,
            last_page=0,
            last_position=legacy_position,
            position_version=None,
        )
        legacy_note = Note(
            title="note",
            content="note content",
            page_number=0,
            position=legacy_position,
            start_pos=None,
            end_pos=None,
            section_title=section.title,
            section_identifier=section.unique_identifier,
            book_id=legacy_book.id,
            position_version=None,
        )
        session.add_all((current_position_info, legacy_position_info, legacy_note))
        session.flush()
        session.execute(
            text("UPDATE book SET content_hash_version = NULL WHERE id = :id"),
            {"id": legacy_book.id},
        )
        session.execute(
            text(
                "UPDATE document_position_info "
                "SET content_hash_version = NULL, position_version = NULL "
                "WHERE id = :id"
            ),
            {"id": legacy_position_info.id},
        )
        session.execute(
            text("UPDATE note SET position_version = NULL WHERE id = :id"),
            {"id": legacy_note.id},
        )
        session.commit()
        session.expire_all()

        reader.load(moved_uri)

        assert Book.query.count() == 1
        assert DocumentPositionInfo.query.count() == 1
        assert Book.query.one().id == current_book.id
        assert DocumentPositionInfo.query.one().id == current_position_info.id
        assert DocumentPositionInfo.query.one().last_position == storage_position
        assert DocumentPositionInfo.query.one().position_version == CURRENT_POSITION_MODEL_VERSION
        assert Note.query.one().book_id == current_book.id
        assert Note.query.one().position == storage_position
        assert Note.query.one().position_version == CURRENT_POSITION_MODEL_VERSION
    finally:
        document.close()
        reader.unload()


def test_reader_merges_legacy_records_when_current_hash_is_unavailable(
    reader,
    tmp_path,
):
    original_path = tmp_path / "original.html"
    moved_path = tmp_path / "moved.html"
    html_content = """
    <html>
        <head><title>Book</title></head>
        <body><p>before <img src="missing.png" alt="Chart"> after text</p></body>
    </html>
    """
    original_path.write_text(html_content, encoding="utf-8")
    moved_path.write_text(html_content, encoding="utf-8")

    original_uri = DocumentUri.from_filename(original_path)
    moved_uri = DocumentUri.from_filename(moved_path)
    document = create_document(moved_uri)
    try:
        assert document.get_content_hash() is None
        legacy_hash = document.get_legacy_content_hash()
        legacy_position = document.get_legacy_content().index("after text")
        storage_position = document.get_storage_content().index("after text")
        section = document.get_page(0).section

        session = Book.session()
        legacy_book = Book(
            title=document.metadata.title,
            uri=original_uri,
            content_hash=legacy_hash,
            content_hash_version=None,
        )
        session.add(legacy_book)
        session.flush()
        session.add_all(
            (
                DocumentPositionInfo(
                    title=document.metadata.title,
                    uri=original_uri,
                    content_hash=legacy_hash,
                    content_hash_version=None,
                    last_page=0,
                    last_position=legacy_position,
                    position_version=None,
                ),
                Note(
                    title="note",
                    content="note content",
                    page_number=0,
                    position=legacy_position,
                    start_pos=None,
                    end_pos=None,
                    section_title=section.title,
                    section_identifier=section.unique_identifier,
                    book_id=legacy_book.id,
                    position_version=None,
                ),
            )
        )
        session.commit()
        session.execute(text("UPDATE book SET content_hash_version = NULL"))
        session.execute(
            text(
                "UPDATE document_position_info "
                "SET content_hash_version = NULL, position_version = NULL"
            )
        )
        session.execute(text("UPDATE note SET position_version = NULL"))
        session.commit()
        session.expire_all()

        reader.load(moved_uri)

        assert Book.query.count() == 1
        assert DocumentPositionInfo.query.count() == 1
        assert Book.query.one().uri == moved_uri
        assert DocumentPositionInfo.query.one().uri == moved_uri
        assert Book.query.one().content_hash == legacy_hash
        assert Book.query.one().content_hash_version is None
        assert DocumentPositionInfo.query.one().content_hash == legacy_hash
        assert DocumentPositionInfo.query.one().content_hash_version is None
        assert DocumentPositionInfo.query.one().last_position == storage_position
        assert DocumentPositionInfo.query.one().position_version == CURRENT_POSITION_MODEL_VERSION
        assert Note.query.one().book_id == Book.query.one().id
        assert Note.query.one().position == storage_position
        assert Note.query.one().position_version == CURRENT_POSITION_MODEL_VERSION

        next_path = tmp_path / "next.html"
        next_path.write_text(html_content, encoding="utf-8")
        next_uri = DocumentUri.from_filename(next_path)

        reader.unload()
        reader.load(next_uri)

        assert Book.query.count() == 1
        assert DocumentPositionInfo.query.count() == 1
        assert Book.query.one().uri == next_uri
        assert Book.query.one().content_hash == legacy_hash
        assert Book.query.one().content_hash_version is None
        assert DocumentPositionInfo.query.one().uri == next_uri
        assert DocumentPositionInfo.query.one().content_hash == legacy_hash
        assert DocumentPositionInfo.query.one().content_hash_version is None
        assert DocumentPositionInfo.query.one().last_position == storage_position
        assert DocumentPositionInfo.query.one().position_version == CURRENT_POSITION_MODEL_VERSION
    finally:
        document.close()
        reader.unload()


def test_legacy_ranges_ending_before_images_migrate_without_including_images(
    reader, tmp_path
):
    html_path = tmp_path / "image.html"
    html_path.write_text(
        """
        <html>
            <head><title>Book</title></head>
            <body>
                <p>before <img src="pic.png" alt="Chart"> middle
                <img src="pic.png" alt="Chart 2"> after text</p>
            </body>
        </html>
        """,
        encoding="utf-8",
    )
    (tmp_path / "pic.png").write_bytes(b"same-image-bytes")
    uri = DocumentUri.from_filename(html_path)
    document = create_document(uri)
    try:
        legacy_text = document.get_legacy_content()
        storage_text = document.get_storage_content()
        legacy_hash = document.get_legacy_content_hash()
        legacy_before_second_image = legacy_text.index("after text")
        first_image = storage_text.index(TEXT_OBJECT_REPLACEMENT_CHAR)
        storage_before_second_image = storage_text.index(
            TEXT_OBJECT_REPLACEMENT_CHAR,
            first_image + 1,
        )
        second_image_start, _second_image_stop = document.get_document_semantic_structure()[
            SemanticElementType.FIGURE
        ][1]
        section = document.get_page(0).section

        session = Book.session()
        book = Book(
            title=document.metadata.title,
            uri=uri,
            content_hash=legacy_hash,
            content_hash_version=None,
        )
        session.add(book)
        session.flush()
        note = Note(
            title="note",
            content="note content",
            page_number=0,
            position=0,
            start_pos=0,
            end_pos=legacy_before_second_image,
            section_title=section.title,
            section_identifier=section.unique_identifier,
            book_id=book.id,
            position_version=None,
        )
        quote = Quote(
            title="quote",
            content="quote content",
            page_number=0,
            position=0,
            start_pos=0,
            end_pos=legacy_before_second_image,
            section_title=section.title,
            section_identifier=section.unique_identifier,
            book_id=book.id,
            position_version=None,
        )
        session.add_all(
            (
                DocumentPositionInfo(
                    title=document.metadata.title,
                    uri=uri,
                    content_hash=legacy_hash,
                    content_hash_version=None,
                    last_page=0,
                    last_position=0,
                    position_version=None,
                ),
                note,
                quote,
            )
        )
        session.flush()
        session.execute(text("UPDATE book SET content_hash_version = NULL"))
        session.execute(
            text(
                "UPDATE document_position_info "
                "SET content_hash_version = NULL, position_version = NULL"
            )
        )
        for table_name in ("note", "quote"):
            session.execute(text(f"UPDATE {table_name} SET position_version = NULL"))
        session.commit()
        session.expire_all()

        reader.load(uri)

        assert Note.query.one().end_pos == storage_before_second_image
        assert Quote.query.one().end_pos == storage_before_second_image
        assert storage_text[Note.query.one().end_pos] == TEXT_OBJECT_REPLACEMENT_CHAR
        assert storage_text[Quote.query.one().end_pos] == TEXT_OBJECT_REPLACEMENT_CHAR
        assert reader.storage_to_view_range(
            Note.query.one().start_pos,
            Note.query.one().end_pos,
            Note.query.one().page_number,
        ).astuple() == (0, second_image_start)
        assert reader.storage_to_view_range(
            Quote.query.one().start_pos,
            Quote.query.one().end_pos,
            Quote.query.one().page_number,
        ).astuple() == (0, second_image_start)
    finally:
        document.close()
        reader.unload()


def test_content_hash_migration_only_adds_schema(asset, tmp_path):
    db_path = tmp_path / "migration.db"
    db_url = f"sqlite:///{db_path}"
    document_uri = DocumentUri.from_filename(asset("test.md")).to_uri_string()

    command.upgrade(_make_alembic_config(db_url), "52e39c4f7494")

    seed_engine = create_engine(db_url, poolclass=NullPool)
    with seed_engine.begin() as conn:
        conn.execute(
            text("INSERT INTO book (title, uri) VALUES (:title, :uri)"),
            {"title": "book", "uri": document_uri},
        )
        conn.execute(
            text(
                "INSERT INTO document_position_info (title, uri, last_page, last_position) "
                "VALUES (:title, :uri, :last_page, :last_position)"
            ),
            {"title": "position", "uri": document_uri, "last_page": 0, "last_position": 0},
        )
        conn.execute(
            text("INSERT INTO recent_document (title, uri) VALUES (:title, :uri)"),
            {"title": "recent", "uri": document_uri},
        )
        conn.execute(
            text(
                "INSERT INTO pinned_document (title, uri, is_pinned, pinning_order) "
                "VALUES (:title, :uri, :is_pinned, :pinning_order)"
            ),
            {"title": "pinned", "uri": document_uri, "is_pinned": True, "pinning_order": 0},
        )
    seed_engine.dispose()

    migrated_engine = init_database(url=db_url, poolclass=NullPool)
    try:
        with migrated_engine.connect() as conn:
            revision = conn.execute(text("SELECT version_num FROM alembic_version")).scalar_one()
            assert revision == "b743b2dbd3a1"
            for table_name in (
                "book",
                "document_position_info",
                "recent_document",
                "pinned_document",
            ):
                content_hashes = (
                    conn.execute(text(f"SELECT content_hash FROM {table_name}")).scalars().all()
                )
                assert content_hashes == [None]
                hash_versions = (
                    conn.execute(text(f"SELECT content_hash_version FROM {table_name}"))
                    .scalars()
                    .all()
                )
                assert hash_versions == [None]
            position_versions = (
                conn.execute(text("SELECT position_version FROM document_position_info"))
                .scalars()
                .all()
            )
            assert position_versions == [None]
    finally:
        close_all_sessions()
        migrated_engine.dispose()


def test_filename_derived_titles_follow_moved_paths(reader, tmp_path):
    first_path = tmp_path / "first.txt"
    second_path = tmp_path / "second.txt"
    first_path.write_text("same text", encoding="utf-8")
    second_path.write_text("same text", encoding="utf-8")

    first_uri = DocumentUri.from_filename(first_path)
    second_uri = DocumentUri.from_filename(second_path)

    reader.load(first_uri)
    first_book = Book.query.one()
    first_position_info = DocumentPositionInfo.query.one()
    reader.view.set_insertion_point(4)
    reader.unload()

    reader.load(second_uri)

    assert Book.query.count() == 1
    assert DocumentPositionInfo.query.count() == 1
    assert Book.query.one().id == first_book.id
    assert DocumentPositionInfo.query.one().id == first_position_info.id
    assert Book.query.one().uri == second_uri
    assert DocumentPositionInfo.query.one().uri == second_uri
    assert reader.stored_document_info.get_last_position() == (0, 4)
    reader.unload()


def test_duplicate_hash_matches_are_merged_into_active_document(reader, tmp_path):
    first_path = tmp_path / "first.txt"
    second_path = tmp_path / "second.txt"
    first_path.write_text("same text", encoding="utf-8")
    second_path.write_text("same text", encoding="utf-8")

    first_uri = DocumentUri.from_filename(first_path)
    second_uri = DocumentUri.from_filename(second_path)
    content_hash = create_document(first_uri).get_content_hash()

    first_book = Book.get_or_create(
        title="first",
        uri=first_uri,
        content_hash=content_hash,
    )
    second_book = Book.get_or_create(
        title="second",
        uri=second_uri,
        content_hash=content_hash,
    )
    first_position_info = DocumentPositionInfo.get_or_create(
        title="first",
        uri=first_uri,
        content_hash=content_hash,
    )
    first_position_info.save_position(0, 7)
    DocumentPositionInfo.get_or_create(
        title="second",
        uri=second_uri,
        content_hash=content_hash,
    )
    Note.session().add(
        Note(
            title="test",
            content="test note",
            page_number=0,
            position=0,
            section_title="test section",
            section_identifier="test-section",
            book_id=first_book.id,
        )
    )
    Note.session().commit()

    reader.load(second_uri)

    assert Book.query.count() == 1
    assert Book.query.one().id == second_book.id
    assert Note.query.one().book_id == second_book.id
    assert NoteTaker(reader).get_for_page(0).count() == 1
    assert DocumentPositionInfo.query.count() == 1
    assert DocumentPositionInfo.query.one().uri == second_uri
    assert reader.stored_document_info.get_last_position() == (0, 7)
    reader.unload()


def test_paginated_content_hash_preserves_page_boundaries():
    first_document = _make_synthetic_paginated_document("first", "alpha", "beta")
    second_document = _make_synthetic_paginated_document("second", "alphab", "eta")
    try:
        assert first_document.get_content_hash() != second_document.get_content_hash()
    finally:
        first_document.close()
        second_document.close()


def test_paginated_documents_with_same_flattened_text_do_not_merge(reader):
    first_document = _make_synthetic_paginated_document(
        "first", "alpha", "beta", title="Paged Document"
    )
    second_document = _make_synthetic_paginated_document(
        "second", "alphab", "eta", title="Paged Document"
    )

    reader.set_document(first_document)
    reader.go_to_page(1, 3)
    reader.save_current_position()
    reader.unload()

    reader.set_document(second_document)

    assert Book.query.count() == 2
    assert DocumentPositionInfo.query.count() == 2
    assert reader.stored_document_info.get_last_position() == (0, 0)
    reader.unload()


def test_textless_documents_do_not_merge_by_content_hash(reader, tmp_path):
    first_path = tmp_path / "first.pptx"
    second_path = tmp_path / "second.pptx"
    _make_textless_presentation(first_path)
    _make_textless_presentation(second_path)

    first_uri = DocumentUri.from_filename(first_path)
    second_uri = DocumentUri.from_filename(second_path)

    first_doc = create_document(first_uri)
    try:
        assert first_doc.get_content_hash() is None
    finally:
        first_doc.close()

    reader.load(first_uri)
    reader.unload()
    reader.load(second_uri)

    assert Book.query.count() == 2
    assert DocumentPositionInfo.query.count() == 2
    assert {record.uri for record in Book.query.all()} == {first_uri, second_uri}
    reader.unload()


@pytest.mark.usefixtures("engine")
def test_image_only_html_documents_do_not_merge_by_content_hash(tmp_path):
    first_path = tmp_path / "first.html"
    second_path = tmp_path / "second.html"
    first_path.write_text(
        '<html><head><title>First</title></head><body><img src="one.png" alt="One"></body></html>',
        encoding="utf-8",
    )
    second_path.write_text(
        '<html><head><title>Second</title></head><body><img src="two.png" alt="Two"></body></html>',
        encoding="utf-8",
    )
    first_document = create_document(DocumentUri.from_filename(first_path))
    second_document = create_document(DocumentUri.from_filename(second_path))
    try:
        assert first_document.get_content_hash() is None
        assert second_document.get_content_hash() is None
        assert first_document.get_legacy_content_hash() is None
        assert second_document.get_legacy_content_hash() is None
        recents_manager.add_to_recents(first_document)
        recents_manager.add_to_recents(second_document)
    finally:
        first_document.close()
        second_document.close()

    assert RecentDocument.query.count() == 2


@pytest.mark.usefixtures("engine")
def test_recent_documents_follow_path_changes_after_lazy_hash_backfill(asset, tmp_path):
    original_uri = DocumentUri.from_filename(asset("roman.epub"))
    original_doc = create_document(original_uri)

    existing_recent = RecentDocument.get_or_create(
        title=original_doc.metadata.title,
        uri=original_uri,
    )
    assert existing_recent.content_hash is None

    recents_manager.add_to_recents(original_doc)

    updated_recent = RecentDocument.query.one()
    assert updated_recent.id == existing_recent.id
    assert updated_recent.content_hash == original_doc.get_content_hash()

    moved_path = shutil.copy(Path(original_uri.path), tmp_path / "roman.epub")
    moved_doc = create_document(DocumentUri.from_filename(moved_path))
    recents_manager.add_to_recents(moved_doc)

    assert RecentDocument.query.count() == 1
    recent = RecentDocument.query.one()
    assert recent.id == existing_recent.id
    assert recent.uri == moved_doc.uri


@pytest.mark.usefixtures("engine")
def test_recent_documents_merge_existing_uri_and_hash_duplicates(asset, tmp_path):
    original_uri = DocumentUri.from_filename(asset("roman.epub"))
    original_doc = create_document(original_uri)
    content_hash = original_doc.get_content_hash()

    moved_path = shutil.copy(Path(original_uri.path), tmp_path / "roman.epub")
    moved_doc = create_document(DocumentUri.from_filename(moved_path))
    current_recent = RecentDocument.get_or_create(
        title=moved_doc.metadata.title,
        uri=moved_doc.uri,
    )
    RecentDocument.get_or_create(
        title=original_doc.metadata.title,
        uri=original_uri,
        content_hash=content_hash,
    )

    try:
        recents_manager.add_to_recents(moved_doc)
    finally:
        original_doc.close()
        moved_doc.close()

    assert RecentDocument.query.count() == 1
    merged_recent = RecentDocument.query.one()
    assert merged_recent.id == current_recent.id
    assert merged_recent.uri == DocumentUri.from_filename(moved_path)
    assert merged_recent.content_hash == content_hash


@pytest.mark.usefixtures("engine")
def test_recent_and_pinned_documents_merge_legacy_when_current_hash_is_unavailable(
    tmp_path,
):
    original_path = tmp_path / "original.html"
    moved_path = tmp_path / "moved.html"
    html_content = """
    <html>
        <head><title>Book</title></head>
        <body><p>before <img src="missing.png" alt="Chart"> after text</p></body>
    </html>
    """
    original_path.write_text(html_content, encoding="utf-8")
    moved_path.write_text(html_content, encoding="utf-8")

    original_uri = DocumentUri.from_filename(original_path)
    moved_document = create_document(DocumentUri.from_filename(moved_path))
    try:
        assert moved_document.get_content_hash() is None
        legacy_hash = moved_document.get_legacy_content_hash()
        legacy_recent = RecentDocument(
            title=moved_document.metadata.title,
            uri=original_uri,
            content_hash=legacy_hash,
            content_hash_version=None,
        )
        legacy_pinned = PinnedDocument(
            title=moved_document.metadata.title,
            uri=original_uri,
            content_hash=legacy_hash,
            content_hash_version=None,
            is_pinned=True,
            pinning_order=7,
        )
        RecentDocument.session.add_all((legacy_recent, legacy_pinned))
        RecentDocument.session.commit()
        RecentDocument.session.execute(
            text("UPDATE recent_document SET content_hash_version = NULL")
        )
        RecentDocument.session.execute(
            text("UPDATE pinned_document SET content_hash_version = NULL")
        )
        RecentDocument.session.commit()
        RecentDocument.session.expire_all()

        recents_manager.add_to_recents(moved_document)
        assert recents_manager.is_pinned(moved_document) is True

        next_path = tmp_path / "next.html"
        next_path.write_text(html_content, encoding="utf-8")
        next_document = create_document(DocumentUri.from_filename(next_path))
        try:
            assert next_document.get_content_hash() is None
            assert next_document.get_legacy_content_hash() == legacy_hash
            recents_manager.add_to_recents(next_document)
            assert recents_manager.is_pinned(next_document) is True
        finally:
            next_document.close()
    finally:
        moved_document.close()

    assert RecentDocument.query.count() == 1
    assert PinnedDocument.query.count() == 1
    recent = RecentDocument.query.one()
    pinned = PinnedDocument.query.one()
    assert recent.uri == next_document.uri
    assert pinned.uri == next_document.uri
    assert recent.content_hash == legacy_hash
    assert pinned.content_hash == legacy_hash
    assert recent.content_hash_version is None
    assert pinned.content_hash_version is None
    assert pinned.is_pinned is True
    assert pinned.pinning_order == 7


@pytest.mark.usefixtures("engine")
def test_textless_recent_documents_do_not_merge_by_content_hash(tmp_path):
    first_path = tmp_path / "first.pptx"
    second_path = tmp_path / "second.pptx"
    _make_textless_presentation(first_path)
    _make_textless_presentation(second_path)

    first_doc = create_document(DocumentUri.from_filename(first_path))
    second_doc = create_document(DocumentUri.from_filename(second_path))
    try:
        recents_manager.add_to_recents(first_doc)
        recents_manager.add_to_recents(second_doc)
    finally:
        first_doc.close()
        second_doc.close()

    assert RecentDocument.query.count() == 2


def test_loading_existing_uri_uses_stored_hash_without_recomputing(reader, asset, monkeypatch):
    uri = DocumentUri.from_filename(asset("roman.epub"))
    document = create_document(uri)
    content_hash = document.get_content_hash()

    Book.get_or_create(
        title=document.metadata.title,
        uri=uri,
        content_hash=content_hash,
        content_hash_version=CURRENT_CONTENT_HASH_VERSION,
    )
    DocumentPositionInfo.get_or_create(
        title=document.metadata.title,
        uri=uri,
        content_hash=content_hash,
        content_hash_version=CURRENT_CONTENT_HASH_VERSION,
    )

    def fail_get_content_hash():
        raise AssertionError("reader should reuse the stored content hash for URI matches")

    monkeypatch.setattr(document, "get_content_hash", fail_get_content_hash)

    reader.set_document(document)

    assert reader.current_book_record.content_hash == content_hash
    assert reader.stored_document_info.content_hash == content_hash
    reader.unload()


@pytest.mark.usefixtures("engine")
def test_recent_documents_with_existing_uri_use_stored_hash_without_recomputing(asset, monkeypatch):
    uri = DocumentUri.from_filename(asset("roman.epub"))
    document = create_document(uri)
    content_hash = document.get_content_hash()

    existing_recent = RecentDocument.get_or_create(
        title=document.metadata.title,
        uri=uri,
        content_hash=content_hash,
        content_hash_version=CURRENT_CONTENT_HASH_VERSION,
    )

    def fail_get_content_hash():
        raise AssertionError("recents should reuse the stored content hash for URI matches")

    monkeypatch.setattr(document, "get_content_hash", fail_get_content_hash)

    try:
        recents_manager.add_to_recents(document)
    finally:
        document.close()

    assert RecentDocument.query.one().id == existing_recent.id


@pytest.mark.usefixtures("engine")
def test_recent_documents_do_not_merge_across_formats(tmp_path):
    markdown_path = tmp_path / "shared.md"
    text_path = tmp_path / "shared.txt"
    markdown_path.write_text("same text", encoding="utf-8")
    text_path.write_text("same text", encoding="utf-8")

    markdown_doc = create_document(DocumentUri.from_filename(markdown_path))
    text_doc = create_document(DocumentUri.from_filename(text_path))
    try:
        recents_manager.add_to_recents(markdown_doc)
        recents_manager.add_to_recents(text_doc)
    finally:
        markdown_doc.close()
        text_doc.close()

    assert RecentDocument.query.count() == 2
    assert {doc.uri.format for doc in RecentDocument.query.all()} == {"markdown", "txt"}


@pytest.mark.usefixtures("engine")
def test_pinned_documents_do_not_match_across_formats(tmp_path):
    markdown_path = tmp_path / "shared.md"
    text_path = tmp_path / "shared.txt"
    markdown_path.write_text("same text", encoding="utf-8")
    text_path.write_text("same text", encoding="utf-8")

    markdown_doc = create_document(DocumentUri.from_filename(markdown_path))
    text_doc = create_document(DocumentUri.from_filename(text_path))
    try:
        recents_manager.pin(markdown_doc)
        assert PinnedDocument.query.count() == 1
        assert recents_manager.is_pinned(markdown_doc) is True
        assert recents_manager.is_pinned(text_doc) is False
    finally:
        markdown_doc.close()
        text_doc.close()

    assert PinnedDocument.query.count() == 2
    assert {doc.uri.format for doc in PinnedDocument.query.all()} == {"markdown", "txt"}
    assert {doc.uri.format for doc in PinnedDocument.query.filter_by(is_pinned=True).all()} == {
        "markdown"
    }


@pytest.mark.usefixtures("engine")
def test_pinned_documents_follow_path_changes_after_lazy_hash_backfill(asset, tmp_path):
    original_uri = DocumentUri.from_filename(asset("roman.epub"))
    original_doc = create_document(original_uri)

    existing_pinned = PinnedDocument.get_or_create(
        title=original_doc.metadata.title,
        uri=original_uri,
    )
    existing_pinned.pin()
    assert existing_pinned.content_hash is None

    assert recents_manager.is_pinned(original_doc) is True

    updated_pinned = PinnedDocument.query.one()
    assert updated_pinned.id == existing_pinned.id
    assert updated_pinned.content_hash == original_doc.get_content_hash()

    moved_path = shutil.copy(Path(original_uri.path), tmp_path / "roman.epub")
    moved_doc = create_document(DocumentUri.from_filename(moved_path))

    assert recents_manager.is_pinned(moved_doc) is True
    assert PinnedDocument.query.count() == 1
    pinned = PinnedDocument.query.one()
    assert pinned.id == existing_pinned.id
    assert pinned.uri == moved_doc.uri


@pytest.mark.usefixtures("engine")
def test_pinned_documents_merge_existing_uri_and_hash_duplicates(asset, tmp_path):
    original_uri = DocumentUri.from_filename(asset("roman.epub"))
    original_doc = create_document(original_uri)
    content_hash = original_doc.get_content_hash()

    moved_path = shutil.copy(Path(original_uri.path), tmp_path / "roman.epub")
    moved_doc = create_document(DocumentUri.from_filename(moved_path))
    current_pinned = PinnedDocument.get_or_create(
        title=moved_doc.metadata.title,
        uri=moved_doc.uri,
    )
    duplicate_pinned = PinnedDocument.get_or_create(
        title=original_doc.metadata.title,
        uri=original_uri,
        content_hash=content_hash,
    )
    duplicate_pinned.pin()
    duplicate_pinned.pinning_order = 7
    PinnedDocument.session.commit()

    try:
        assert recents_manager.is_pinned(moved_doc) is True
    finally:
        original_doc.close()
        moved_doc.close()

    assert PinnedDocument.query.count() == 1
    merged_pinned = PinnedDocument.query.one()
    assert merged_pinned.id == current_pinned.id
    assert merged_pinned.uri == DocumentUri.from_filename(moved_path)
    assert merged_pinned.content_hash == content_hash
    assert merged_pinned.is_pinned is True
    assert merged_pinned.pinning_order == 7


@pytest.mark.usefixtures("engine")
def test_pinned_documents_merge_legacy_duplicate_with_existing_uri(asset, tmp_path):
    original_uri = DocumentUri.from_filename(asset("roman.epub"))
    moved_path = shutil.copy(Path(original_uri.path), tmp_path / "roman.epub")
    moved_document = create_document(DocumentUri.from_filename(moved_path))
    try:
        moved_uri = moved_document.uri
        current_hash = moved_document.get_content_hash()
        legacy_hash = moved_document.get_legacy_content_hash()
        current_pinned = PinnedDocument.get_or_create(
            title=moved_document.metadata.title,
            uri=moved_uri,
        )
        current_pinned.content_hash = None
        current_pinned.content_hash_version = None
        current_pinned.is_pinned = False
        current_pinned.pinning_order = 0
        legacy_pinned = PinnedDocument(
            title=moved_document.metadata.title,
            uri=original_uri,
            content_hash=legacy_hash,
            content_hash_version=None,
            is_pinned=True,
            pinning_order=7,
        )
        PinnedDocument.session.add(legacy_pinned)
        PinnedDocument.session.flush()
        PinnedDocument.session.execute(
            text(
                "UPDATE pinned_document "
                "SET content_hash = NULL, content_hash_version = NULL, is_pinned = 0, pinning_order = 0 "
                "WHERE id = :id"
            ),
            {"id": current_pinned.id},
        )
        PinnedDocument.session.execute(
            text(
                "UPDATE pinned_document "
                "SET content_hash_version = NULL, is_pinned = 1, pinning_order = 7 "
                "WHERE id = :id"
            ),
            {"id": legacy_pinned.id},
        )
        PinnedDocument.session.commit()
        PinnedDocument.session.expire_all()

        assert recents_manager.is_pinned(moved_document) is True
    finally:
        moved_document.close()

    assert PinnedDocument.query.count() == 1
    merged_pinned = PinnedDocument.query.one()
    assert merged_pinned.id == current_pinned.id
    assert merged_pinned.uri == moved_uri
    assert merged_pinned.content_hash == current_hash
    assert merged_pinned.content_hash_version == CURRENT_CONTENT_HASH_VERSION
    assert merged_pinned.is_pinned is True
    assert merged_pinned.pinning_order == 7
